from pptx import Presentation
from pptx.util import Inches
from app.api.logger import setup_logger
import os

logger = setup_logger(__name__)

def return_images():
    images = [
        [
            {
                'path': f'{os.getcwd()}/app/api/features/images/Python-Symbol.png',
                'left': 5,
                'top': 4,
                'width': 2.5,
                'height': 2
            }
        ],
        [
            {
                'path': f'{os.getcwd()}/app/api/features/images/code.jpg',
                'left': 5,
                'top': 4,
                'width': 2.5,
                'height': 2
            },
        ]
    ]

    return images

# Load the template presentation
template_path = f'{os.getcwd()}/app/api/features/templates/template1.pptx'

def create_pptx_file(result, images):
    prs = Presentation(template_path)

    # Remove all existing slides
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    file_title = result['title'].replace(" ", "_")

    # Add a title slide
    slide_layout = prs.slide_layouts[0]  # Assuming the first layout is the title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = result['title']
    subtitle.text = result['description']

    # Add content slides
    for slide_index, slide_content in enumerate(result['slides']):
        slide_layout = prs.slide_layouts[1]  # Assuming the second layout is a content slide layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content_shape = slide.placeholders[1]

        title.text = slide_content['title']
        content_shape.text = slide_content['content']

        # Add images to the slide if any
        if slide_index < len(images):
            for image in images[slide_index]:
                left = Inches(image['left'])
                top = Inches(image['top'])
                width = Inches(image['width'])
                height = Inches(image['height'])
                slide.shapes.add_picture(image['path'], left, top, width=width, height=height)

    # Save the presentation
    logger.info("Creating new PPT file")
    pptx_file = f"{os.getcwd()}/app/api/features/results/{file_title}.pptx"
    prs.save(pptx_file)
    logger.info("The PPTX file is saved successfully")